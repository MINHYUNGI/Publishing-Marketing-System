from __future__ import annotations
from pathlib import Path
from typing import Iterable
from .config import MAX_DOCUMENT_TEXT

def _clean(lines: Iterable[str]) -> str:
    result = []
    for line in lines:
        value = " ".join(str(line).replace("\x00", " ").split())
        if value:
            result.append(value)
    return "\n".join(result)[:MAX_DOCUMENT_TEXT]

def read_document_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".xlsx":
        return _read_xlsx(path)
    if suffix == ".xls":
        return _read_xls(path)
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    raise ValueError(f"지원하지 않는 문서 형식입니다: {suffix}")

def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True, read_only=True)
    lines = [f"[파일명] {path.name}"]
    try:
        for sheet in wb.worksheets:
            lines.append(f"[시트] {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if values:
                    lines.append(" | ".join(values))
    finally:
        wb.close()
    return _clean(lines)

def _read_xls(path: Path) -> str:
    import xlrd
    book = xlrd.open_workbook(path)
    lines = [f"[파일명] {path.name}"]
    for sheet in book.sheets():
        lines.append(f"[시트] {sheet.name}")
        for r in range(sheet.nrows):
            values = [str(v).strip() for v in sheet.row_values(r) if str(v).strip()]
            if values:
                lines.append(" | ".join(values))
    return _clean(lines)

def _read_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    lines = [f"[파일명] {path.name}"]
    for i, page in enumerate(reader.pages, 1):
        lines.append(f"[페이지 {i}]")
        lines.append(page.extract_text() or "")
    return _clean(lines)

def _read_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    lines = [f"[파일명] {path.name}"]
    lines.extend(p.text for p in doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            lines.append(" | ".join(cell.text for cell in row.cells))
    return _clean(lines)

def _read_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = [f"[파일명] {path.name}"]
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[슬라이드 {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return _clean(lines)


def _number(value):
    if value is None or value == "":
        return None
    if isinstance(value, (int, float)):
        return int(round(value))
    text = str(value).strip().replace(",", "").replace("원", "").replace("부", "")
    try:
        return int(round(float(text)))
    except Exception:
        return None


def _norm(value) -> str:
    return "".join(str(value or "").replace("\n", " ").split()).lower()


def _extract_sales_targets_from_rows(rows: list[list[object]]) -> dict:
    """표준 마케팅 기획서의 영업목표 표를 셀 위치 기반으로 읽습니다.

    AI 추정이 아니라 원본 Excel 셀의 숫자를 그대로 보존합니다.
    헤더명(초도, 3/6/12개월, BEP)과 행명(부수, 매출액)을 찾아 교차 셀을 읽습니다.
    """
    if not rows:
        return {}

    header_aliases = {
        "initial": ("초도", "초도배본"),
        "month3": ("3개월", "~3개월", "출간~3개월"),
        "month6": ("6개월", "~6개월", "출간~6개월"),
        "month12": ("12개월", "~12개월", "출간~12개월"),
        "bep": ("bep",),
    }
    best = None
    for ri, row in enumerate(rows):
        normalized = [_norm(v) for v in row]
        cols = {}
        for key, aliases in header_aliases.items():
            for ci, text in enumerate(normalized):
                if any(alias.replace(" ", "").lower() in text for alias in aliases):
                    cols[key] = ci
                    break
        score = len(cols)
        if score >= 3 and (best is None or score > best[0]):
            best = (score, ri, cols)
    if not best:
        return {}

    _, header_row, cols = best
    units_row = sales_row = None
    for ri in range(max(0, header_row - 3), min(len(rows), header_row + 12)):
        texts = [_norm(v) for v in rows[ri]]
        joined = "|".join(texts)
        if units_row is None and ("부수" in joined or "매출부수" in joined):
            units_row = ri
        if sales_row is None and "매출액" in joined:
            sales_row = ri
    if units_row is None and sales_row is None:
        return {}

    out = {}
    key_map = {
        "initial": ("initial_units", "initial_sales"),
        "month3": ("month3_units", "month3_sales"),
        "month6": ("month6_units", "month6_sales"),
        "month12": ("month12_units", "month12_sales"),
        "bep": ("bep_units", "bep_sales"),
    }
    for period, ci in cols.items():
        uk, sk = key_map[period]
        if units_row is not None and ci < len(rows[units_row]):
            out[uk] = _number(rows[units_row][ci])
        if sales_row is not None and ci < len(rows[sales_row]):
            out[sk] = _number(rows[sales_row][ci])

    # BEP 초과 목표 시점은 별도 셀에서 보존(있을 때만).
    for ri, row in enumerate(rows):
        for ci, value in enumerate(row):
            text = _norm(value)
            if "bep" in text and ("초과" in text or "시점" in text):
                candidates = []
                if ci + 1 < len(row): candidates.append(row[ci + 1])
                if ri + 1 < len(rows) and ci < len(rows[ri + 1]): candidates.append(rows[ri + 1][ci])
                for candidate in candidates:
                    ctext = str(candidate or "").strip()
                    if ctext:
                        out["bep_target_note"] = ctext
                        import re
                        m = re.search(r"(\d+)\s*개월", ctext)
                        if m: out["bep_target_month"] = int(m.group(1))
                        break
                if out.get("bep_target_note"):
                    break
        if out.get("bep_target_note"):
            break

    return {k: v for k, v in out.items() if v is not None and v != ""}


def extract_sales_targets(path: Path) -> dict:
    """Excel 마케팅 기획서의 영업 목표 숫자를 직접 추출합니다."""
    suffix = path.suffix.lower()
    try:
        if suffix == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(path, data_only=True, read_only=True)
            try:
                best = {}
                for sheet in wb.worksheets:
                    rows = [list(row) for row in sheet.iter_rows(values_only=True)]
                    found = _extract_sales_targets_from_rows(rows)
                    if len(found) > len(best): best = found
                return best
            finally:
                wb.close()
        if suffix == ".xls":
            import xlrd
            book = xlrd.open_workbook(path)
            best = {}
            for sheet in book.sheets():
                rows = [sheet.row_values(r) for r in range(sheet.nrows)]
                found = _extract_sales_targets_from_rows(rows)
                if len(found) > len(best): best = found
            return best
    except Exception:
        return {}
    return {}
